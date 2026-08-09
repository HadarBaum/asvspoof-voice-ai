"""
Builds slides.pptx for the course presentation. Reads real numbers from
docs/training_metrics.json when available (falls back to a "run the
pipeline first" placeholder so this is safe to run at any point).

Usage (from the repo root):
    python docs/slides/generate_slides.py
"""

import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt

DOCS_DIR = os.path.join(os.path.dirname(__file__), "..")
METRICS_PATH = os.path.join(DOCS_DIR, "training_metrics.json")
CHARTS_DIR = os.path.join(DOCS_DIR, "..", "app", "static", "charts")
OUT_PATH = os.path.join(os.path.dirname(__file__), "slides.pptx")


def load_metrics():
    if os.path.exists(METRICS_PATH):
        with open(METRICS_PATH) as f:
            return json.load(f)
    return None


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
    return slide


def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
    return slide


def main():
    metrics = load_metrics()
    prs = Presentation()

    add_title_slide(
        prs,
        "Human-or-AI Voice Detector",
        "BIU Big Data and AI course project — ASVspoof2019 (LA)",
    )

    add_bullets_slide(
        prs,
        "Problem & dataset",
        [
            "Can we tell real human speech apart from AI-generated (TTS / voice-conversion) speech?",
            "ASVspoof2019 Logical Access partition: bonafide vs. spoofed utterances,",
            "~121k labeled flac clips across train / dev / eval, speaker-disjoint splits.",
            "PA (replay-attack) partition intentionally excluded — not AI-generated speech.",
        ],
    )

    add_bullets_slide(
        prs,
        "Architecture",
        [
            "MinIO (object store) holds the raw audio corpus",
            "Spark (batch, parallel) extracts acoustic features from every clip",
            "Elasticsearch (NoSQL) stores features, training results, and live predictions",
            "Kafka streams simulated incoming audio events",
            "A Kafka consumer scores each event in near-real-time with the trained model",
            "Flask app: /classify (live demo) and /dashboard (insights)",
        ],
    )

    add_bullets_slide(
        prs,
        "AI capability",
        [
            "Random Forest classifier on acoustic summary features (MFCC, spectral, pitch, energy)",
            "Trained on the dataset's own speaker-disjoint train/dev split",
            "Applied both offline (batch validation) and online (streaming enrichment)",
            "Same model backs the live web demo — one implementation, not three",
        ],
    )

    if metrics:
        add_bullets_slide(
            prs,
            "Results",
            [
                f"Trained on {metrics['n_train']} utterances, validated on {metrics['n_dev']}",
                f"Dev accuracy: {metrics['accuracy_dev']:.1%}",
                f"Dev precision: {metrics['precision_dev']:.1%} | recall: {metrics['recall_dev']:.1%} | F1: {metrics['f1_dev']:.1%}",
            ],
        )
    else:
        add_bullets_slide(
            prs,
            "Results",
            ["Run the pipeline (see README.md) and re-run this script to fill in real numbers."],
        )

    accuracy_chart = os.path.join(CHARTS_DIR, "accuracy_by_attack.png")
    if os.path.exists(accuracy_chart):
        add_image_slide(prs, "Detection accuracy by attack type", accuracy_chart)

    add_bullets_slide(
        prs,
        "Challenges & trade-offs",
        [
            "Windows + Spark: worker Python subprocess resolution breaks under paths with spaces",
            "Chose a plain Kafka consumer over Spark Structured Streaming for the enrichment step",
            "  (avoids the spark-sql-kafka connector for a comparatively light per-message workload)",
            "Kept Python end-to-end (not Scala) since the feature-extraction library is Python-only",
            "Speaker-disjoint train/dev split (not random) to avoid inflating accuracy via leakage",
        ],
    )

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
